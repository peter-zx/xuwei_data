import os
import uuid
import shutil
import subprocess
import platform
from pathlib import Path
from typing import List, Optional, Callable
from datetime import datetime
import time

from config.settings import get_settings
from app.core.logging import get_logger
from app.models.schemas import ConversionResultData


class ConversionResult:
    def __init__(
        self,
        success: bool,
        original_path: str,
        output_path: str = "",
        error: str = "",
        file_size: int = 0,
        processing_time: float = 0.0
    ):
        self.success = success
        self.original_path = original_path
        self.output_path = output_path
        self.error = error
        self.file_size = file_size
        self.processing_time = processing_time
    
    def to_data(self) -> ConversionResultData:
        return ConversionResultData(
            success=self.success,
            original_path=self.original_path,
            output_path=self.output_path,
            error=self.error,
            file_size=self.file_size,
            processing_time=self.processing_time
        )


class Doc2PdfConverter:
    def __init__(self, output_dir: str, source_root: str = None):
        self.settings = get_settings()
        self.logger = get_logger("converter")
        
        if output_dir:
            self.output_dir = Path(output_dir)
        else:
            self.output_dir = Path.home() / "Desktop" / f"Output_{datetime.now().strftime('%m%d_%H%M%S')}"
        
        self.output_dir.mkdir(parents=True, exist_ok=True)
        self.source_root = Path(source_root) if source_root else None
    
    def convert(self, file_path: str, progress_callback: Optional[Callable] = None) -> ConversionResult:
        start_time = time.time()
        path = Path(file_path)
        
        if not path.exists():
            return ConversionResult(
                success=False,
                original_path=file_path,
                error=f"文件不存在: {file_path}"
            )
        
        ext = path.suffix.lower()
        
        if ext not in self.settings.conversion.SUPPORTED_EXTENSIONS:
            return ConversionResult(
                success=False,
                original_path=file_path,
                error=f"不支持的格式: {ext}"
            )
        
        try:
            output_path = self._get_output_path(path)
            
            if ext == '.pdf':
                shutil.copy(file_path, output_path)
                processing_time = time.time() - start_time
                return ConversionResult(
                    success=True,
                    original_path=file_path,
                    output_path=str(output_path),
                    file_size=output_path.stat().st_size,
                    processing_time=processing_time
                )
            
            if ext in ['.docx']:
                result = self._convert_word(file_path, output_path)
            elif ext == '.doc':
                result = self._convert_doc_with_com(file_path, output_path)
            elif ext in ['.xlsx', '.xls']:
                result = self._convert_excel(file_path, output_path)
            elif ext in ['.pptx', '.ppt']:
                result = self._convert_ppt(file_path, output_path)
            elif ext == '.txt':
                result = self._convert_txt(file_path, output_path)
            else:
                result = ConversionResult(
                    success=False,
                    original_path=file_path,
                    error=f"不支持的格式: {ext}"
                )
            
            if result.success:
                result.processing_time = time.time() - start_time
                result.file_size = output_path.stat().st_size if output_path.exists() else 0
            
            if progress_callback:
                progress_callback()
            
            return result
            
        except Exception as e:
            self.logger.exception(f"转换失败: {file_path}")
            return ConversionResult(
                success=False,
                original_path=file_path,
                error=str(e)
            )
    
    def _get_output_path(self, path: Path) -> Path:
        if self.source_root:
            source = Path(self.source_root)
            rel_path = path.relative_to(source)
            output_subdir = self.output_dir / source.name / rel_path.parent
            output_subdir.mkdir(parents=True, exist_ok=True)
            output_name = f"{path.stem}.pdf"
            return output_subdir / output_name
        else:
            output_name = f"{path.stem}_{uuid.uuid4().hex[:8]}.pdf"
            return self.output_dir / output_name
    
    def _convert_word(self, input_path: str, output_path: Path) -> ConversionResult:
        system = platform.system()
        
        if system == "Windows":
            result = self._convert_with_wps(input_path, output_path)
            if result.success:
                return result
        
        return self._convert_with_libreoffice(input_path, output_path)
    
    def _convert_with_libreoffice(self, input_path: str, output_path: Path) -> ConversionResult:
        libreoffice_path = shutil.which('libreoffice') or shutil.which('soffice')
        
        if not libreoffice_path:
            return ConversionResult(
                success=False,
                original_path=input_path,
                error="请安装LibreOffice进行转换"
            )
        
        try:
            input_dir = str(Path(input_path).parent)
            
            cmd = [
                libreoffice_path,
                '--headless',
                '--convert-to', 'pdf',
                '--outdir', input_dir,
                input_path
            ]
            
            result = subprocess.run(cmd, capture_output=True, text=True, timeout=120)
            
            generated_pdf = Path(input_dir) / f"{Path(input_path).stem}.pdf"
            if generated_pdf.exists():
                shutil.move(str(generated_pdf), str(output_path))
                return ConversionResult(success=True, original_path=input_path, output_path=str(output_path))
            else:
                return ConversionResult(success=False, original_path=input_path, error="LibreOffice转换失败")
                
        except subprocess.TimeoutExpired:
            return ConversionResult(success=False, original_path=input_path, error="转换超时")
        except Exception as e:
            return ConversionResult(success=False, original_path=input_path, error=str(e))
    
    def _convert_doc_with_com(self, input_path: str, output_path: Path) -> ConversionResult:
        system = platform.system()
        
        if system == "Windows":
            result = self._convert_with_wps(input_path, output_path)
            if result.success:
                return result
        
        return self._convert_with_libreoffice(input_path, output_path)
    
    def _convert_with_wps(self, input_path: str, output_path: Path) -> ConversionResult:
        try:
            import pythoncom
            import win32com.client
        except ImportError:
            return ConversionResult(success=False, original_path=input_path, error="WPS/Office仅Windows支持，云端请使用LibreOffice")
        
        try:
            
            success = False
            for prog_id in ["WPS.Application", "KSO.Application", "Word.Application"]:
                try:
                    try:
                        pythoncom.CoInitialize()
                    except Exception:
                        pass
                    
                    app = win32com.client.Dispatch(prog_id)
                    doc = app.Documents.Open(str(Path(input_path).absolute()), ReadOnly=True)
                    doc.SaveAs(str(output_path), FileFormat=17)
                    doc.Close()
                    app.Quit()
                    pythoncom.CoUninitialize()
                    
                    if output_path.exists():
                        success = True
                        break
                except Exception as e:
                    try:
                        pythoncom.CoUninitialize()
                    except:
                        pass
                    continue
            
            if success:
                return ConversionResult(success=True, original_path=input_path, output_path=str(output_path))
            else:
                return ConversionResult(success=False, original_path=input_path, error="WPS/Office转换失败")
            
        except Exception as e:
            self.logger.error(f"WPS conversion error: {e}")
            return ConversionResult(success=False, original_path=input_path, error=f"WPS转换异常: {str(e)}")
    
    def _convert_excel(self, input_path: str, output_path: Path) -> ConversionResult:
        try:
            from openpyxl import load_workbook
            from reportlab.lib.pagesizes import A4, landscape
            from reportlab.lib.styles import getSampleStyleSheet
            from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Table
            from reportlab.lib.units import cm
            from reportlab.lib import colors
            
            wb = load_workbook(input_path, data_only=True)
            
            pdf_doc = SimpleDocTemplate(
                str(output_path),
                pagesize=landscape(A4),
                leftMargin=1*cm,
                rightMargin=1*cm,
                topMargin=1*cm,
                bottomMargin=1*cm
            )
            styles = getSampleStyleSheet()
            story = []
            
            for sheet_name in wb.sheetnames:
                ws = wb[sheet_name]
                story.append(Paragraph(sheet_name, styles['Heading2']))
                
                table_data = []
                max_cols = min(ws.max_column, 10)
                
                for row in ws.iter_rows(max_row=min(ws.max_row, 50), max_col=max_cols):
                    row_data = [str(cell.value) if cell.value else '' for cell in row[:max_cols]]
                    table_data.append(row_data)
                
                if table_data:
                    t = Table(table_data, repeatRows=1)
                    t.setStyle([
                        ('BACKGROUND', (0, 0), (-1, 0), colors.lightgrey),
                        ('TEXTCOLOR', (0, 0), (-1, 0), colors.black),
                        ('FONTNAME', (0, 0), (-1, -1), 'Helvetica'),
                        ('FONTSIZE', (0, 0), (-1, -1), 8),
                        ('ALIGN', (0, 0), (-1, -1), 'CENTER'),
                        ('GRID', (0, 0), (-1, -1), 0.5, colors.grey),
                        ('VALIGN', (0, 0), (-1, -1), 'MIDDLE'),
                    ])
                    story.append(t)
                story.append(Spacer(1, 20))
            
            pdf_doc.build(story)
            return ConversionResult(success=True, original_path=input_path, output_path=str(output_path))
            
        except Exception as e:
            return ConversionResult(success=False, original_path=input_path, error=f"Excel转换失败: {str(e)}")
    
    def _convert_ppt(self, input_path: str, output_path: Path) -> ConversionResult:
        try:
            from pptx import Presentation
            from reportlab.lib.pagesizes import A4
            from reportlab.lib.styles import getSampleStyleSheet
            from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer
            from reportlab.lib.units import cm
            
            prs = Presentation(input_path)
            
            pdf_doc = SimpleDocTemplate(
                str(output_path),
                pagesize=A4,
                leftMargin=2*cm,
                rightMargin=2*cm,
                topMargin=2*cm,
                bottomMargin=2*cm
            )
            styles = getSampleStyleSheet()
            story = []
            
            for i, slide in enumerate(prs.slides, 1):
                story.append(Paragraph(f"Slide {i}", styles['Heading3']))
                
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        for para in shape.text.split('\n'):
                            if para.strip():
                                story.append(Paragraph(para, styles['Normal']))
                
                story.append(Spacer(1, 20))
            
            pdf_doc.build(story)
            return ConversionResult(success=True, original_path=input_path, output_path=str(output_path))
            
        except Exception as e:
            return ConversionResult(success=False, original_path=input_path, error=f"PPT转换失败: {str(e)}")
    
    def _convert_txt(self, input_path: str, output_path: Path) -> ConversionResult:
        try:
            with open(input_path, 'r', encoding='utf-8') as f:
                content = f.read()
            
            from reportlab.lib.pagesizes import A4
            from reportlab.lib.styles import getSampleStyleSheet
            from reportlab.platypus import SimpleDocTemplate, Paragraph
            from reportlab.lib.units import cm
            
            pdf_doc = SimpleDocTemplate(
                str(output_path),
                pagesize=A4,
                leftMargin=2*cm,
                rightMargin=2*cm,
                topMargin=2*cm,
                bottomMargin=2*cm
            )
            styles = getSampleStyleSheet()
            story = []
            
            for line in content.split('\n'):
                story.append(Paragraph(line or ' ', styles['Normal']))
            
            pdf_doc.build(story)
            return ConversionResult(success=True, original_path=input_path, output_path=str(output_path))
            
        except Exception as e:
            return ConversionResult(success=False, original_path=input_path, error=f"TXT转换失败: {str(e)}")
    
    def convert_batch(
        self,
        file_paths: List[str],
        progress_callback: Optional[Callable[[int, int], None]] = None
    ) -> List[ConversionResult]:
        total = len(file_paths)
        completed = 0
        results = []
        
        self.logger.info(f"Starting batch conversion: {total} files")
        
        for path in file_paths:
            result = self.convert(path)
            results.append(result)
            completed += 1
            
            if progress_callback:
                progress_callback(completed, total)
            
            status = "SUCCESS" if result.success else "FAILED"
            self.logger.info(f"[{completed}/{total}] {status}: {path}")
        
        successful = sum(1 for r in results if r.success)
        self.logger.info(f"Batch conversion completed: {successful}/{total} successful")
        
        return results


def get_converter(output_dir: Optional[str] = None, source_root: Optional[str] = None) -> Doc2PdfConverter:
    return Doc2PdfConverter(output_dir, source_root)
