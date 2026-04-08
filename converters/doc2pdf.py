import os
import subprocess
from pathlib import Path
from typing import List, Dict, Optional, Callable
import tempfile
import shutil


class ConversionResult:
    def __init__(self, success: bool, original_path: str, output_path: str = "", error: str = ""):
        self.success = success
        self.original_path = original_path
        self.output_path = output_path
        self.error = error


class Doc2PdfConverter:
    def __init__(self, output_dir: str, source_root: str = None):
        self.output_dir = Path(output_dir)
        self.output_dir.mkdir(parents=True, exist_ok=True)
        self.source_root = Path(source_root) if source_root else None
        self.results: List[ConversionResult] = []
    
    def convert_single(self, file_path: str, progress_callback: Optional[Callable] = None) -> ConversionResult:
        path = Path(file_path)
        ext = path.suffix.lower()
        
        try:
            if self.source_root:
                rel_path = Path(file_path).relative_to(Path(self.source_root))
                output_subdir = self.output_dir / rel_path.parent
                output_subdir.mkdir(parents=True, exist_ok=True)
                output_path = output_subdir / f"{path.stem}.pdf"
            else:
                output_path = self.output_dir / f"{path.stem}.pdf"
            
            if ext == '.pdf':
                shutil.copy(file_path, output_path)
                return ConversionResult(True, file_path, str(output_path))
            
            if ext in ['.docx']:
                result = self._convert_word(file_path, output_path)
            elif ext == '.doc':
                result = self._convert_doc_old(file_path, output_path)
            elif ext in ['.xlsx', '.xls']:
                result = self._convert_excel(file_path, output_path)
            elif ext in ['.pptx', '.ppt']:
                result = self._convert_ppt(file_path, output_path)
            elif ext == '.txt':
                result = self._convert_txt(file_path, output_path)
            else:
                return ConversionResult(False, file_path, error=f"Unsupported format: {ext}")
            
            if progress_callback:
                progress_callback()
            
            return result
            
        except Exception as e:
            return ConversionResult(False, file_path, error=str(e))
    
    def _convert_word(self, input_path: str, output_path: Path) -> ConversionResult:
        try:
            from docx import Document
            from reportlab.lib.pagesizes import A4
            from reportlab.lib.styles import getSampleStyleSheet
            from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Table
            from reportlab.lib.units import cm
            
            doc = Document(input_path)
            
            pdf_doc = SimpleDocTemplate(str(output_path), pagesize=A4, 
                                        leftMargin=2*cm, rightMargin=2*cm,
                                        topMargin=2*cm, bottomMargin=2*cm)
            styles = getSampleStyleSheet()
            story = []
            
            for para in doc.paragraphs:
                text = para.text.strip()
                if text:
                    if para.style.name.startswith('Heading'):
                        level = int(para.style.name[-1]) if para.style.name[-1].isdigit() else 1
                        story.append(Paragraph(text, styles[f'Heading{level}']))
                    else:
                        story.append(Paragraph(text, styles['Normal']))
                    story.append(Spacer(1, 6))
            
            for table in doc.tables:
                table_data = []
                for row in table.rows:
                    row_data = [cell.text for cell in row.cells]
                    table_data.append(row_data)
                if table_data:
                    t = Table(table_data)
                    story.append(t)
                    story.append(Spacer(1, 12))
            
            pdf_doc.build(story)
            return ConversionResult(True, input_path, str(output_path))
            
        except Exception as e:
            return ConversionResult(False, input_path, error=f"Word conversion failed: {str(e)}")
    
    def _convert_doc_old(self, input_path: str, output_path: Path) -> ConversionResult:
        import subprocess
        import shutil
        import platform
        
        system = platform.system()
        
        if system == "Windows":
            result = self._convert_with_wps(input_path, output_path)
            if result.success:
                return result
        
        libreoffice_path = shutil.which('libreoffice') or shutil.which('soffice')
        
        if not libreoffice_path:
            return ConversionResult(
                False, 
                input_path, 
                error=".doc not supported by python-docx. Save as .docx to convert."
            )
        
        try:
            input_dir = str(Path(input_path).parent)
            
            cmd = [
                libreoffice_path,
                '--headless',
                '--convert-to', 'pdf',
                '--outdir', input_dir,
                input_path
            ]
            
            result = subprocess.run(cmd, capture_output=True, text=True, timeout=60)
            
            generated_pdf = Path(input_dir) / f"{Path(input_path).stem}.pdf"
            if generated_pdf.exists():
                generated_pdf.rename(output_path)
                return ConversionResult(True, input_path, str(output_path))
            else:
                return ConversionResult(False, input_path, error="LibreOffice conversion failed")
                
        except subprocess.TimeoutExpired:
            return ConversionResult(False, input_path, error="Conversion timeout")
        except Exception as e:
            return ConversionResult(False, input_path, error=str(e))
    
    def _convert_with_wps(self, input_path: str, output_path: Path) -> ConversionResult:
        import pythoncom
        import win32com.client
        
        for prog_id in ["WPS.Application", "KSO.Application", "Word.Application"]:
            try:
                pythoncom.CoInitialize()
                
                app = win32com.client.Dispatch(prog_id)
                doc = app.Documents.Open(str(Path(input_path).absolute()), ReadOnly=True)
                doc.SaveAs(str(output_path), FileFormat=17)
                doc.Close()
                app.Quit()
                
                pythoncom.CoUninitialize()
                
                if output_path.exists():
                    return ConversionResult(True, input_path, str(output_path))
                    
            except Exception as e:
                try:
                    pythoncom.CoUninitialize()
                except:
                    pass
                continue
        
        return ConversionResult(False, input_path, error=".doc转换失败，请安装WPS或Office")
    
    def _convert_excel(self, input_path: str, output_path: Path) -> ConversionResult:
        try:
            from openpyxl import load_workbook
            from reportlab.lib.pagesizes import A4, landscape
            from reportlab.lib import colors
            from reportlab.lib.styles import getSampleStyleSheet
            from reportlab.platypus import SimpleDocTemplate, Table, Spacer
            from reportlab.lib.units import cm
            
            wb = load_workbook(input_path, data_only=True)
            
            pdf_doc = SimpleDocTemplate(str(output_path), pagesize=landscape(A4),
                                        leftMargin=1*cm, rightMargin=1*cm,
                                        topMargin=1*cm, bottomMargin=1*cm)
            styles = getSampleStyleSheet()
            story = []
            
            for sheet_name in wb.sheetnames:
                ws = wb[sheet_name]
                story.append(Paragraph(sheet_name, styles['Heading2']))
                
                table_data = []
                max_cols = min(ws.max_column, 10)
                
                for row in ws.iter_rows(max_row=min(ws.max_row, 50), max_col=max_cols):
                    row_data = [str(cell.value) if cell.value else '' for cell in row[:max_cols]]
                    table_data.append(row_data)
                
                if table_data:
                    t = Table(table_data, repeatRows=1)
                    t.setStyle([
                        ('BACKGROUND', (0, 0), (-1, 0), colors.lightgrey),
                        ('TEXTCOLOR', (0, 0), (-1, 0), colors.black),
                        ('FONTNAME', (0, 0), (-1, -1), 'Helvetica'),
                        ('FONTSIZE', (0, 0), (-1, -1), 8),
                        ('ALIGN', (0, 0), (-1, -1), 'CENTER'),
                        ('GRID', (0, 0), (-1, -1), 0.5, colors.grey),
                        ('VALIGN', (0, 0), (-1, -1), 'MIDDLE'),
                    ])
                    story.append(t)
                story.append(Spacer(1, 20))
            
            pdf_doc.build(story)
            return ConversionResult(True, input_path, str(output_path))
            
        except Exception as e:
            return ConversionResult(False, input_path, error=f"Excel conversion failed: {str(e)}")
    
    def _convert_ppt(self, input_path: str, output_path: Path) -> ConversionResult:
        try:
            from pptx import Presentation
            from reportlab.lib.pagesizes import A4
            from reportlab.lib.styles import getSampleStyleSheet
            from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer
            from reportlab.lib.units import cm
            
            prs = Presentation(input_path)
            
            pdf_doc = SimpleDocTemplate(str(output_path), pagesize=A4,
                                        leftMargin=2*cm, rightMargin=2*cm,
                                        topMargin=2*cm, bottomMargin=2*cm)
            styles = getSampleStyleSheet()
            story = []
            
            for i, slide in enumerate(prs.slides, 1):
                story.append(Paragraph(f"Slide {i}", styles['Heading3']))
                
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        for para in shape.text.split('\n'):
                            if para.strip():
                                story.append(Paragraph(para, styles['Normal']))
                
                story.append(Spacer(1, 20))
            
            pdf_doc.build(story)
            return ConversionResult(True, input_path, str(output_path))
            
        except Exception as e:
            return ConversionResult(False, input_path, error=f"PPT conversion failed: {str(e)}")
    
    def _convert_txt(self, input_path: str, output_path: Path) -> ConversionResult:
        try:
            with open(input_path, 'r', encoding='utf-8') as f:
                content = f.read()
            
            from reportlab.lib.pagesizes import A4
            from reportlab.lib.styles import getSampleStyleSheet
            from reportlab.platypus import SimpleDocTemplate, Paragraph
            from reportlab.lib.units import cm
            
            pdf_doc = SimpleDocTemplate(str(output_path), pagesize=A4,
                                        leftMargin=2*cm, rightMargin=2*cm,
                                        topMargin=2*cm, bottomMargin=2*cm)
            styles = getSampleStyleSheet()
            story = []
            
            for line in content.split('\n'):
                story.append(Paragraph(line or ' ', styles['Normal']))
            
            pdf_doc.build(story)
            return ConversionResult(True, input_path, str(output_path))
            
        except Exception as e:
            return ConversionResult(False, input_path, error=f"TXT conversion failed: {str(e)}")
    
    def convert_batch(self, file_paths: List[str], 
                     progress_callback: Optional[Callable[[int, int], None]] = None) -> List[ConversionResult]:
        total = len(file_paths)
        completed = 0
        results = []
        
        for path in file_paths:
            result = self.convert_single(path)
            results.append(result)
            completed += 1
            if progress_callback:
                progress_callback(completed, total)
        
        return results
